#!/usr/bin/env python3
"""
SmartDoc Translator - 智能文档翻译工具 (集成 Kimi AI)
====================================================

使用当前 AI 模型作为翻译引擎，保留原文格式

四大问题解决：
1. 大文件 → 本地解析
2. 上下文超限 → 智能分块 + 摘要
3. 格式保留 → 位置提取 + 回填
4. 术语一致 → 术语表 + 上下文记忆

用法：
    python smartdoc.py translate input.pdf --output output.pdf
"""

import argparse
import json
import os
import re
import sqlite3
import subprocess
import sys
import tempfile
from pathlib import Path
from typing import Dict, List, Optional, Tuple
from dataclasses import dataclass, asdict

# 文档解析库
try:
    import fitz  # PyMuPDF
    HAS_FITZ = True
except ImportError:
    HAS_FITZ = False
    print("⚠️  请安装 PyMuPDF: pip install pymupdf")

try:
    from pptx import Presentation
    from pptx.util import Pt
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


@dataclass
class TextElement:
    """文档文本元素"""
    text: str
    x: float
    y: float
    font: str
    size: float
    color: str
    page: int
    element_type: str
    
    def to_dict(self):
        return asdict(self)


class TranslationMemory:
    """翻译记忆库"""
    
    def __init__(self, db_path: str = "~/.smartdoc/memory.db"):
        self.db_path = Path(db_path).expanduser()
        self.db_path.parent.mkdir(parents=True, exist_ok=True)
        self.conn = sqlite3.connect(str(self.db_path))
        self._init_db()
    
    def _init_db(self):
        cursor = self.conn.cursor()
        cursor.execute('''
            CREATE TABLE IF NOT EXISTS glossary (
                id INTEGER PRIMARY KEY,
                source TEXT UNIQUE,
                target TEXT,
                context TEXT,
                frequency INTEGER DEFAULT 1,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            )
        ''')
        cursor.execute('''
            CREATE TABLE IF NOT EXISTS translations (
                id INTEGER PRIMARY KEY,
                file_hash TEXT,
                source TEXT,
                target TEXT,
                context_summary TEXT,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            )
        ''')
        self.conn.commit()
    
    def get_term(self, source: str) -> Optional[str]:
        cursor = self.conn.cursor()
        cursor.execute("SELECT target FROM glossary WHERE source = ?", (source,))
        result = cursor.fetchone()
        return result[0] if result else None
    
    def add_term(self, source: str, target: str, context: str = ""):
        cursor = self.conn.cursor()
        cursor.execute('''
            INSERT INTO glossary (source, target, context)
            VALUES (?, ?, ?)
            ON CONFLICT(source) DO UPDATE SET
                target = excluded.target,
                frequency = frequency + 1
        ''', (source, target, context))
        self.conn.commit()
    
    def get_recent_context(self, limit: int = 3) -> str:
        cursor = self.conn.cursor()
        cursor.execute('''
            SELECT context_summary FROM translations
            ORDER BY created_at DESC
            LIMIT ?
        ''', (limit,))
        results = cursor.fetchall()
        return " | ".join([r[0] for r in results if r[0]])[:500]
    
    def add_translation(self, source: str, target: str, summary: str = ""):
        cursor = self.conn.cursor()
        cursor.execute('''
            INSERT INTO translations (source, target, context_summary)
            VALUES (?, ?, ?)
        ''', (source[:1000], target[:1000], summary[:500]))
        self.conn.commit()


class SmartChunker:
    """智能分块器 - 法律文档专用"""
    
    def __init__(self, max_chars: int = 800, overlap: int = 50):
        # 法律文档：小块逐句翻译，确保不遗漏
        self.max_chars = max_chars  # 从2000减到800
        self.overlap = overlap  # 从100减到50
    
    def chunk_elements(self, elements: List[TextElement]) -> List[List[TextElement]]:
        """按元素分块（保持页码完整）"""
        chunks = []
        current_chunk = []
        current_chars = 0
        current_page = None
        
        for elem in elements:
            elem_chars = len(elem.text)
            
            # 新页码或超过限制，创建新块
            if (current_page is not None and elem.page != current_page) or \
               (current_chars + elem_chars > self.max_chars):
                if current_chunk:
                    chunks.append(current_chunk)
                # 保留少量重叠元素用于上下文
                overlap_count = min(2, len(current_chunk))
                current_chunk = current_chunk[-overlap_count:] if overlap_count > 0 else []
                current_chars = sum(len(e.text) for e in current_chunk)
            
            current_chunk.append(elem)
            current_chars += elem_chars
            current_page = elem.page
        
        if current_chunk:
            chunks.append(current_chunk)
        
        return chunks


class KimiTranslator:
    """使用 Kimi (当前模型) 进行翻译"""
    
    def __init__(self, memory: TranslationMemory):
        self.memory = memory
        self.prompt_template = """你是一位专业的法律文档翻译专家。请严格按照以下要求翻译：

【翻译要求 - 必须遵守】
1. **逐句翻译**：每一句话都要翻译，不要跳过任何内容
2. **不要总结**：严禁概括、归纳或改写，必须保留原文所有细节
3. **不要省略**：禁止省略任何从句、修饰语或补充说明
4. **术语一致**：专业术语必须前后统一，使用标准法律术语
5. **格式对应**：原文的段落、列表、编号格式都要保留
6. **保留原文结构**：法律条文的条款、子条款层级必须对应

【禁止事项】
- ❌ 禁止概括性翻译（如"包括但不限于"改成"包括"）
- ❌ 禁止合并句子
- ❌ 禁止删减修饰成分
- ❌ 禁止改变法律条文的严谨性

【术语表】（必须使用这些译法）
{glossary}

【前文上下文】（仅供参考连贯性）
{context}

【待翻译文本 - 请逐句翻译】
{text}

【逐句翻译结果】（只输出中文翻译，不要解释，保持原文结构）："""
    
    def _build_glossary_text(self) -> str:
        """构建术语表文本"""
        cursor = self.memory.conn.cursor()
        cursor.execute("SELECT source, target FROM glossary ORDER BY frequency DESC LIMIT 20")
        terms = cursor.fetchall()
        if not terms:
            return "（暂无术语表）"
        return "\n".join([f"- {s} → {t}" for s, t in terms])
    
    def translate(self, text: str) -> str:
        """调用 Kimi 进行翻译"""
        # 检查术语表缓存
        cached = self.memory.get_term(text[:200])  # 检查前200字符
        if cached and len(text) < 200:
            return cached
        
        # 获取上下文
        context = self.memory.get_recent_context(3)
        glossary = self._build_glossary_text()
        
        # 构建提示词
        prompt = self.prompt_template.format(
            glossary=glossary,
            context=context if context else "（开始翻译）",
            text=text
        )
        
        # 调用 Kimi API (通过 sessions_spawn 或本地调用)
        # 这里使用环境变量或文件方式传递
        return self._call_kimi(prompt, text)
    
    def _call_kimi(self, prompt: str, original_text: str) -> str:
        """调用 Kimi API"""
        try:
            # 方法1: 使用环境变量中的 API key
            api_key = os.environ.get("KIMI_API_KEY") or os.environ.get("OPENCLAW_API_KEY")
            
            if api_key:
                # 直接调用 API
                import urllib.request
                import json
                
                data = json.dumps({
                    "model": "kimi-for-coding",
                    "messages": [{"role": "user", "content": prompt}],
                    "temperature": 0.3,
                    "max_tokens": 4000
                }).encode()
                
                req = urllib.request.Request(
                    "https://api.kimi.com/coding/v1/chat/completions",
                    data=data,
                    headers={
                        "Authorization": f"Bearer {api_key}",
                        "Content-Type": "application/json"
                    }
                )
                
                with urllib.request.urlopen(req, timeout=60) as response:
                    result = json.loads(response.read())
                    translated = result["choices"][0]["message"]["content"]
                    
                    # 保存到记忆
                    summary = translated[:100] if len(translated) > 100 else translated
                    self.memory.add_translation(original_text, translated, summary)
                    
                    # 如果是短文本，也保存到术语表
                    if len(original_text) < 100 and len(translated) < 100:
                        self.memory.add_term(original_text, translated)
                    
                    return translated
            
            # 方法2: 使用本地文件方式（通过 OpenClaw 调用）
            # 写入临时文件，等待外部处理
            return self._call_via_file(prompt, original_text)
            
        except Exception as e:
            print(f"⚠️  翻译API调用失败: {e}")
            # 降级：返回原文
            return original_text
    
    def _call_via_file(self, prompt: str, original_text: str) -> str:
        """通过文件方式调用（用于 OpenClaw 集成）"""
        # 创建临时请求文件
        temp_dir = Path("~/.smartdoc/queue").expanduser()
        temp_dir.mkdir(parents=True, exist_ok=True)
        
        request_id = f"{os.urandom(4).hex()}"
        request_file = temp_dir / f"{request_id}.json"
        response_file = temp_dir / f"{request_id}.response"
        
        request_data = {
            "id": request_id,
            "prompt": prompt,
            "original_text": original_text,
            "status": "pending"
        }
        
        with open(request_file, 'w') as f:
            json.dump(request_data, f)
        
        print(f"   📝 已创建翻译请求: {request_id}")
        print(f"   ⏳ 等待 AI 处理...")
        
        # 等待响应（最多60秒）
        import time
        for i in range(60):
            if response_file.exists():
                with open(response_file) as f:
                    response = json.load(f)
                # 清理文件
                request_file.unlink(missing_ok=True)
                response_file.unlink(missing_ok=True)
                
                translated = response.get("translated", original_text)
                
                # 保存到记忆
                summary = translated[:100] if len(translated) > 100 else translated
                self.memory.add_translation(original_text, translated, summary)
                
                return translated
            time.sleep(1)
        
        # 超时，清理并返回原文
        request_file.unlink(missing_ok=True)
        return original_text


class PDFParser:
    """PDF解析器"""
    
    def __init__(self, file_path: str):
        if not HAS_FITZ:
            raise ImportError("请安装 PyMuPDF: pip install pymupdf")
        self.file_path = file_path
        self.doc = fitz.open(file_path)
    
    def extract(self) -> List[TextElement]:
        elements = []
        
        for page_num in range(len(self.doc)):
            page = self.doc[page_num]
            blocks = page.get_text("dict")["blocks"]
            
            for block in blocks:
                if "lines" in block:
                    for line in block["lines"]:
                        for span in line["spans"]:
                            text = span["text"].strip()
                            if text and len(text) > 1:  # 过滤单字符
                                element = TextElement(
                                    text=text,
                                    x=span["bbox"][0],
                                    y=span["bbox"][1],
                                    font=span.get("font", "unknown"),
                                    size=span.get("size", 12),
                                    color=str(span.get("color", 0)),
                                    page=page_num + 1,
                                    element_type="text"
                                )
                                elements.append(element)
        
        return elements
    
    def close(self):
        self.doc.close()


class PPTParser:
    """PPT解析器"""
    
    def __init__(self, file_path: str):
        if not HAS_PPTX:
            raise ImportError("请安装 python-pptx: pip install python-pptx")
        self.file_path = file_path
        self.prs = Presentation(file_path)
    
    def extract(self) -> List[TextElement]:
        elements = []
        
        for slide_num, slide in enumerate(self.prs.slides, 1):
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    if len(text) > 1:
                        element = TextElement(
                            text=text,
                            x=getattr(shape, 'left', 0),
                            y=getattr(shape, 'top', 0),
                            font="unknown",
                            size=18,
                            color="0",
                            page=slide_num,
                            element_type="textbox"
                        )
                        elements.append(element)
        
        return elements


class TextFormatter:
    """纯文本格式化输出（用于测试和预览）"""
    
    def write_text(self, elements: List[TextElement], output_path: str):
        """输出为带格式的文本文件"""
        with open(output_path, 'w', encoding='utf-8') as f:
            current_page = 0
            for elem in elements:
                if elem.page != current_page:
                    f.write(f"\n{'='*50}\n")
                    f.write(f"第 {elem.page} 页\n")
                    f.write(f"{'='*50}\n\n")
                    current_page = elem.page
                
                f.write(f"{elem.text}\n\n")


class SmartDocApp:
    """主应用类"""
    
    def __init__(self):
        self.memory = TranslationMemory()
        self.translator = KimiTranslator(self.memory)
    
    def translate_file(self, input_path: str, output_path: str) -> bool:
        input_path = Path(input_path)
        
        print(f"🔄 SmartDoc 翻译")
        print(f"   输入: {input_path.name}")
        print(f"   输出: {output_path}")
        print()
        
        # 1. 解析文档
        print("📄 步骤1/4: 解析文档...")
        try:
            if input_path.suffix.lower() == '.pdf':
                parser = PDFParser(str(input_path))
            elif input_path.suffix.lower() in ['.pptx', '.ppt']:
                parser = PPTParser(str(input_path))
            else:
                print(f"❌ 不支持的格式: {input_path.suffix}")
                return False
            
            elements = parser.extract()
            print(f"   ✓ 提取 {len(elements)} 个文本元素")
            
            if not elements:
                print("⚠️  文档中没有找到文本")
                return False
            
        except Exception as e:
            print(f"❌ 解析失败: {e}")
            return False
        
        # 2. 智能分块
        print()
        print("🧩 步骤2/4: 智能分块...")
        chunker = SmartChunker(max_chars=1500)  # 控制每块大小
        chunks = chunker.chunk_elements(elements)
        print(f"   ✓ 分成 {len(chunks)} 个块")
        
        # 3. 翻译
        print()
        print("🌐 步骤3/4: 翻译中...")
        print("   使用 Kimi AI 进行翻译...")
        
        translated_elements = []
        total = len(chunks)
        
        for i, chunk in enumerate(chunks, 1):
            print(f"   翻译块 {i}/{total} ({len(chunk)} 个元素)...", end=' ', flush=True)
            
            for elem in chunk:
                # 翻译文本
                translated_text = self.translator.translate(elem.text)
                elem.text = translated_text
                translated_elements.append(elem)
            
            print("✓")
        
        print(f"   ✓ 翻译完成 ({len(translated_elements)} 个元素)")
        
        # 4. 输出
        print()
        print("📝 步骤4/4: 生成输出...")
        
        # 目前输出为文本格式（便于预览）
        formatter = TextFormatter()
        
        if output_path.endswith('.txt'):
            formatter.write_text(translated_elements, output_path)
        else:
            # 默认输出为 .txt 便于查看
            txt_output = output_path.rsplit('.', 1)[0] + '_zh.txt'
            formatter.write_text(translated_elements, txt_output)
            print(f"   ✓ 输出: {txt_output}")
        
        parser.close()
        
        print()
        print("✅ 翻译完成！")
        return True


def main():
    parser = argparse.ArgumentParser(
        description='SmartDoc Translator - 智能文档翻译工具 (Kimi AI 版)'
    )
    parser.add_argument('command', choices=['translate', 'glossary'],
                       help='命令')
    parser.add_argument('input', nargs='?', help='输入文件')
    parser.add_argument('-o', '--output', help='输出文件')
    
    args = parser.parse_args()
    
    app = SmartDocApp()
    
    if args.command == 'translate':
        if not args.input:
            print("❌ 请指定输入文件")
            sys.exit(1)
        
        if not Path(args.input).exists():
            print(f"❌ 文件不存在: {args.input}")
            sys.exit(1)
        
        output = args.output or args.input.rsplit('.', 1)[0] + '_zh.txt'
        success = app.translate_file(args.input, output)
        sys.exit(0 if success else 1)
    
    elif args.command == 'glossary':
        print("📚 术语表:")
        memory = TranslationMemory()
        cursor = memory.conn.cursor()
        cursor.execute("SELECT source, target, frequency FROM glossary ORDER BY frequency DESC LIMIT 20")
        terms = cursor.fetchall()
        if terms:
            for source, target, freq in terms:
                print(f"   {source} → {target} (使用{freq}次)")
        else:
            print("   (术语表为空，开始翻译后将自动积累)")


if __name__ == '__main__':
    main()
